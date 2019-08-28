#import json
#from urllib import parse
import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import plotly.graph_objs as go
import plotly.io as pio
#import numpy as np
import plotly
import matplotlib.pyplot as plt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.text import PP_ALIGN
import calendar
from selenium import webdriver
#from pyxlsb import open_workbook as open_xlsb
from PIL import Image
from urllib.request import urlopen
from io import BytesIO
from webdriver_manager.chrome import ChromeDriverManager
import sys
#-----------------------------------------------------------------------------#
sys.path.insert(0, os.getcwd())
#directory = os.path.dirname(__file__) + '/'
directory = '/tmp/'
#print(directory)
ppt_directory = 'file://' + directory
image_dir = ''
pio.orca.config.executable = '/user_code/' + 'plotly-orca-1.2.1-1/orca_app/orca.exe'
#mapbox_access_token = 'pk.eyJ1IjoiYW5hbHl0aWNzbG4iLCJhIjoiY2p0eWFrbzJ4MGZ6czRkcG5tc3hka3A3MiJ9.3iT-wbkITehGa-nhf5AgTw'
#io.orca.config.mapbox_access_token = mapbox_access_token
pio.orca.config.save()

class v:
    Sno = 'index'
    raw_loc_U = 'Raw Location'
    raw_state_U = 'Raw State'
    raw_loc = 'Raw Location_S'
    raw_state = 'Raw State_S'
    clean_loc = 'Clean Location_S'
    clean_state = 'Clean State_S'
    Admin1 = 'Districtname'
    Admin2 = 'Sub-distname'
    Admin3 = 'Officename'
    Admin4 = 'Village/Locality name'
    SuperAdmin = 'StateName'
    typee = 'type'
    PIN = 'Pincode'
    pinlist = 'pinlist'
    key_remove = ['\'', '\xa0', 'Taluka', 'City', 'Ex', 'Dist', 'Distt',
                  'District', 'Store', 'Agri', 'Godown', 'Plant', 'Factory',
                  'Fresh', 'Bulk', 'Kasba', 'Bazaar', 'Bazar', 'Works', 'Work',
                  'Tql', 'Tq.', 'Rly.', 'Rly', 'Del @', 'Dl At', 'Del ', 'Port',
                  'Point', 'Cwh', 'Rdc']
    
    user_input = 'user_input'
    tool_sugg = 'Suggestion'
    manual_io = 'manual'
    opt_rej = 'Rejected'
    opt_othr = 'Other'
    opt_acpt = 'Accepted'
    
    # ----------------Mapping Module Variable --------------#
    """ Variable for any mapping is fix and sheet heading specially for database
     should be same """
    
    dbkey = 'Raw'
    dbvalue = 'Map'
    vkey = 'Vendor_S'
    tkey = 'Truck_S'
    
    # ------------ For integration ---------------------#
    rawO = 'raworiginname'
    rawOstate = 'raworiginstate'
    rawD = 'rawdestinationname'
    rawDstate = 'rawdestinationstate'
    rawV = 'rawvendor'
    rawTrk = 'rawtrucktype'
    network_id = 'dataset'
    country = 'country'
    sector = 'sector'
    timestamp = 'receiveddate'
    rate = 'rate'
    vol = 'volume'
    
    # ---------------added by data cleaning app ---------- #
    Opoint = 'OriginPoint'
    Odis = 'OriginDistrict'
    Ostate = 'OriginState'
    OPIN = 'OriginPincode'
    Dpoint = 'DestinationPoint'
    Ddis = 'DestinationDistrict'
    Dstate = 'DestinationState'
    DPIN = 'DestinationPincode'
    C_ven = 'Vendor'
    C_trucktype = 'Trucktype'
    ven_id = 'transporter_id'
    ven_name = 'transporter_name'

from azure.storage.blob import BlockBlobService

def connInitate():
    account_key = ('2e/1qLktFuzffwkhbyC2++FwUgMzAigJqiNHxghNINDx25Tg/MggLUy6Se0R'
                   + 'Vb4jGiIpFsIZJc2INpRKEQTO+g==')
    block_blob_service = BlockBlobService('logisticsnowtech3', account_key)
    return block_blob_service

def writeBlob(file_path, filename, container_name='rajkumar'):
    block_blob_service = connInitate()
    block_blob_service.create_blob_from_path(container_name, r'{}'.format(filename),
                                             r'{}'.format(file_path))
    return print("Thanks")

'''Functions Used'''
#-----------------------------------------------------------------------------#
def run_text(slide_name,l,t,w,h,subtitle,font_size,r,g,b):
    left = Inches(l)
    top = Inches(t)
    width = Inches(w)
    height = Inches(h)
    txBox = slide_name.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = subtitle
    font = run.font
    font.color.rgb = RGBColor(r, g, b)
    font.size = Pt(font_size)

def move_slide(presentation, old_index, new_index):
    """Used for changing Slide Positions.
    
    :param prs: Presentation instance of python-pptx
    :type prs: pptx.presentation.Presentation
    :param old_index: Index of slide where it is residing
    :type old_index: int
    :param new_index: Index of slide where you want to move
    :type new_index: int
    """
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])
    
def summarySlide(df):
    """Function used for calculationg the summary of Datarame
    containing ::

        Total Spend (in lac)
        Total Volume
        # Lanes
        # Origin
        # Destination
        # Unique Truck
        # Unique Vendor
   
    :param df: Dataframe of which summary has to be made
    :type df: DataFrame
    :return: DataFrame with Summary
    :rtype: DataFrame
    """
    if {'spend'}.issubset(df.columns):
        df['spend'] = df['spend'].fillna(-1)
        if (df['spend'].sum() <= 0) & ({'rate', 'volume'}.issubset(df.columns)):
            df['rate'] = pd.to_numeric(df['rate'], errors='coerce')
            df['volume'] = pd.to_numeric(df['volume'], errors='coerce')
            df['spend'] = df['rate'] * df['volume']
        else:
            pass
    elif {'rate', 'volume'}.issubset(df.columns):
        df['rate'] = pd.to_numeric(df['rate'], errors='coerce')
        df['volume'] = pd.to_numeric(df['volume'], errors='coerce')
        df['spend'] = df['rate'] * df['volume']
    else:
        pass
    unique_lane = len(pd.unique(df[v.Opoint] + df[v.Ostate]
                                + df[v.Dpoint] + df[v.Dstate]
                                + df[v.C_trucktype]))
    unique_origin = len(pd.unique(df[v.Opoint] + df[v.Ostate]))
    unique_destination = len(pd.unique(df[v.Dpoint] + df[v.Dstate]))
    unique_truck = len(pd.unique(df[v.C_trucktype]))
    unique_vendr = len(pd.unique(df[v.C_ven]))
    total_spend = df['spend'].sum() / 1e5
    total_vol = df[v.vol].sum()
    data = [[total_spend, total_vol, unique_lane,
            unique_origin, unique_destination, unique_truck,
            unique_vendr]]
    col = ['total spend', 'total vol', '# lane', '# origin',
           '# destination', '# unique_truck', '# unique_vendr']
    df = pd.DataFrame(data=data, columns=col)
    return df

def top(df, col_name, based_on, upto_number, number_format=1,
        percentage=False):
    if {'spend'}.issubset(df.columns):
        pass
    elif {'rate', 'volume'}.issubset(df.columns):
        df['rate'] = pd.to_numeric(df['rate'], errors='coerce')
        df['volume'] = pd.to_numeric(df['volume'], errors='coerce')
        df['spend'] = df['rate'] * df['volume']
    else:
        pass
    if percentage:
        summ = df[based_on].sum()/number_format
    else:
        pass
    if isinstance(based_on, str):
        df = df.groupby([col_name]).agg({based_on: sum})
        df = df.reset_index()
        df[based_on] = round(df[based_on]/number_format, 2)
        df = df.sort_values(based_on, ascending=False)
        df = df.head(n=upto_number)
        df = df.reset_index(drop=True)
        for idx, item in enumerate(df[col_name]):
            # print(idx, item)
            if len(item) > 20:
                df.loc[idx, col_name] = item[:20]
    if percentage:
        df[based_on] = df[based_on]/summ
    return df

def barChartConstruct(prs, typee, title, df, based_on, axis_title,
                      color='b'):
    """Used for Making Bar chart
    
    :param prs: Presentation instance of python-pptx
    :type prs: pptx.presentation.Presentation
    :param typee: Used to decide Columns of Categorial Variable
    :type typee: str
    :param title: Title of Presentation Slide
    :type title: str
    :param df: Dataset of which PPT is in construction
    :type df: DataFrame
    :param based_on: Continuous variable on which Calculation has to performed
    :type based_on: str
    :param axis_title: Title of Axis of Categorial Variable
    :type axis_title: str
    :param color: If given 'b' ppt bars will be in blue colur, defaults to 'b'
    :type color: str, optional
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    Title_placeholder = slide.placeholders[0]
    run_text(slide,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
    Title_placeholder.text = title
#    Title_placeholder.text_frame.margin_top = Inches()
    chart_data = CategoryChartData()
    chart_data.categories = df[typee].iloc[::-1]
    chart_data.add_series('Series 1', df[based_on].iloc[::-1])
    x, y, cx, cy = Inches(1.25), Inches(1.75), Inches(7), Inches(5)
    chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.minor_tick_mark = XL_TICK_MARK.INSIDE
#    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(15)
    category_axis.tick_labels.font.bold = False
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = max(df[based_on]) + 0.05 * max(df[based_on])
    value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_minor_gridlines = False
    value_axis.has_major_gridlines = False
    value_axis.axis_title.text_frame.text = axis_title
    tick_labels = value_axis.tick_labels
#    tick_labels.font.bold = True
    tick_labels.font.size = Pt(14)
    plot = chart.plots[0]
    plot.vary_by_categories = False
    series = plot.series[0]
    fill = series.format.fill
    fill.solid()
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(13)
    if color == 'b':
        data_labels.font.color.rgb = RGBColor(0,100,0)
    else:
        fill.fore_color.rgb = RGBColor(0,100,0)
        data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    return slide

def getBlob(blob_name_incont, download_path, container_name='rajkumar'):
    block_blob_service = connInitate()
    block_blob_service.get_blob_to_path(container_name, blob_name_incont,
                                        download_path)

def htmlToimg(file):
    #downloadpath = directory+image_dir+"chromedriver.exe"
    #getBlob("chromedriver.exe", downloadpath)
#    downloadpath = "https://logisticsnowtech3.blob.core.windows.net/rajkumar/chromedriver.exe"
    driver = webdriver.Chrome('/user_code/' + "chromedriver.exe")
#    driver = webdriver.Chrome(ChromeDriverManager().install())
    driver.get(ppt_directory +file)
    save_name = 'map.png'
    import time  
    time.sleep(7)    
    driver.save_screenshot(directory+image_dir+ save_name)
    driver.quit()
    img = Image.open(directory+image_dir+ save_name)
    box = (10, 10, 600, 380)
    area = img.crop(box)
    area.save(directory+image_dir+ save_name)
    return None
#---------------------------------map---------------------------------#
def Lane_map(prs, df, origin, dest):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title4 = slide.placeholders[0]
    title4.text = "Network of Lanes"
    title4.text_frame.margin_top = Inches(-0.5)
    mapbox_access_token = 'pk.eyJ1IjoiYW5hbHl0aWNzbG4iLCJhIjoiY2p0eWFrbzJ4MGZ6czRkcG5tc3hka3A3MiJ9.3iT-wbkITehGa-nhf5AgTw'
    style='mapbox://styles/analyticsln/cjv66wfoy0jv71fs3188uxl6x'
    dfh=df.drop_duplicates([origin, dest])
    center=go.layout.mapbox.Center(lat=21.153,lon=79.083)
    zoom=3.3
    layout = go.Layout(autosize=False,width = 600,height = 380, hovermode='closest',showlegend=True,
               legend = dict(x= 0.015, y= 0.98,font = dict(size = 10)),margin={'l': 1, 'b': 1, 't': 1,'r': 1},
                       mapbox=dict(accesstoken=mapbox_access_token,bearing=0,
                       pitch=0, zoom=zoom, style=style, center=center
                       ))
    data=[]
    for i in range(0,len(dfh)):
        data.append(
                    go.Scattermapbox(lon = [dfh['OriginLongitude'].iloc[i], dfh['DestinationLongitude'].iloc[i]],
        lat = [dfh['OriginLatitude'].iloc[i], dfh['DestinationLatitude'].iloc[i]],
        mode = 'lines',showlegend=False,
        line = dict(width =0.5,color = 'turquoise'),))

    lon1=pd.DataFrame(dfh['OriginLongitude'])
    lon2=pd.DataFrame(dfh['DestinationLongitude'])

    lat1=pd.DataFrame(dfh['OriginLatitude'])
    lat2=pd.DataFrame(dfh['DestinationLatitude'])

    name1=pd.DataFrame(dfh[origin])
    name2=pd.DataFrame(dfh[dest])
    point1=[go.Scattermapbox(lat = lat1['OriginLatitude'], lon=lon1['OriginLongitude'], mode='markers+text', marker=dict(size=8,color='red'),
                             text=name1[origin],textposition='top center',showlegend=True,
                             textfont=dict(family='sans serif',size=10,color='white'),
                             name='Origin')]
    point2=[go.Scattermapbox(lat = lat2['DestinationLatitude'], lon=lon2['DestinationLongitude'], mode='markers+text', marker=dict(size=5,color='blue'),
                             text=name2[dest],textposition='top center',showlegend=True,
                             textfont=dict(family='sans serif',size=10,color='white'),
                             name='Destination')]
    map_fig = {'data': data+point1+point2, 'layout': layout}
    pio.write_image(map_fig, directory+image_dir+'map.png')
    plotly.offline.plot(map_fig, filename=directory+image_dir+"map.html")
    #htmlToimg(image_dir+"map.html")
    vendor_image1 = directory+image_dir+'map.png'
    slide.shapes.add_picture(vendor_image1,Inches(0.8),Inches(1.5))
    run_text(slide,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
    return None

#------------------------summary by volume-----------------------------#

def summary_chart(prs, df, date, based_on):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    title.text = "Summary – Volume by Month"
    df[date] = pd.to_datetime(df[date])
    df['months'] = df[date].dt.month
    df1 = df.groupby(['months']).agg({based_on: sum})
    df1 = df1.sort_values('months')
    df1 = df1.reset_index()
    df1['months'] = df1['months'].apply(lambda x: calendar.month_abbr[x])
    plt.plot(list(df1['months']), list(df1[based_on]), 'C3', label='Months',
             zorder=1, lw=2)
    plt.scatter(list(df1['months']), list(df1[based_on]),s=50, zorder=2)
    plt.ylabel("Volume (in Truckload)", fontsize=15,fontweight='medium')
    plt.xlabel("Months", fontsize=15,fontweight='medium')
    plt.tight_layout()
    plt.savefig(directory+image_dir+'fig2.png')
    vendor_image = directory+image_dir+'fig2.png'
    slide.shapes.add_picture(vendor_image,Inches(1.7),Inches(1.7))    
    run_text(slide,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
    return None

#-----------------------bar + line -------------------------------------------#

def bar_lineChart(prs, df, col, based_on1, based_on2):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    if col == v.Opoint:
        title.text = "Top 10 Origins by Spend and Volume"
    elif col == v.Dpoint:
        title.text = "Top 10 Destinations by Spend and Volume"
    elif col == 'Lane':
        title.text = "Top 10 Lanes by Spend and Volume"
    title.size = Pt(20)
    df_spend = top(df, col, based_on1, 10, number_format=1e5)
    df_volume = top(df, col, based_on2, 50)
    df1 = df_spend.merge(df_volume, left_on=col, right_on=col)
    fig, ax1 = plt.subplots()
    ax1.set_xlabel(col, fontsize=15,fontweight='medium')
    ax1.set_ylabel("Spend (in Lakh)", fontsize=15,fontweight='medium')
    color = 'tab:blue'
    plt.bar(list(df1[col]), list(df1[based_on1]), color=color, label='Spend', width=0.4)
    ax1.tick_params(axis='y')
    ax1.legend(loc='upper right', bbox_to_anchor=(1, 0.9))
    ax2 = ax1.twinx()
    color = 'tab:red'
    ax2.set_ylabel("Volume (in Truckload)", fontsize=15,fontweight='medium')
    ax2.plot(list(df1[col]), list(df1[based_on2]), color=color, label='Volume')
    ax2.tick_params(axis='y')
    ax2.legend(loc='upper right', bbox_to_anchor=(1, 1))
    if len(df_spend) > 6:
        fig.set_figheight(5)
        fig.set_figwidth(10)
        plt.xlabel(r'$x$')
    else:
        fig.set_figheight(5)
        fig.set_figwidth(8)        
    plt.savefig(directory+image_dir+'fig3.png')
    vendor_image5 = directory+image_dir+'fig3.png'
    if len(df_spend) > 6:
        slide.shapes.add_picture(vendor_image5,Inches(0),Inches(1.5))
    else:
        slide.shapes.add_picture(vendor_image5,Inches(1),Inches(1.5))
    run_text(slide,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
    return None 



#-------------------------donut chart-----------------------------------------#
def donutChart(prs, df, col1, based_on):
    def func(data):
        perc = []
        for i in data:
            perc.append(str(round(i/np.sum(data)*100,1))+'%')
        return perc
    slide_layout = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(slide_layout)
    title3 = slide3.placeholders[0]
    if based_on == 'spend':
        title3.text = "{} By Spend".format(col1)
    elif based_on == 'volume':
        title3.text = "{} By Volume".format(col1)
    df_spend = top(df, col1, based_on, 10)
    fig1, ax1 = plt.subplots()
    explode = tuple([0.05]*len(df_spend))
    data = list(df_spend[based_on])
    label=list(df_spend[col1])
    patches, texts = ax1.pie(data, labels=func(data), startangle=90,
                             explode = explode)
    centre_circle = plt.Circle((0,0),0.7,fc='white')
    fig = plt.gcf()
    fig.gca().add_artist(centre_circle)
    plt.legend(label,loc='upper right')
    plt.margins(0.1,0.1)
    fig.set_size_inches(8, 5)
    ax1.axis('equal')  
    plt.tight_layout()
    plt.savefig(directory+image_dir+'fig4.png')
    vendor_image1 = directory+image_dir+'fig4.png'
    slide3.shapes.add_picture(vendor_image1,Inches(0.8),Inches(1.5))
    run_text(slide3,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
    return None


#-----------------Stacked bar chart------------------------------------#
def stackedBar(prs, df, col1, col2, value):
    slide_layout = prs.slide_layouts[5]
    slide1 = prs.slides.add_slide(slide_layout)
    title3 = slide1.placeholders[0]
    title3.text = "Volume Distribution by Trucktype"
    title3.text_frame.margin_top = Inches(-0.5)
    table = pd.pivot_table(df, values=value, index=[col1, col2],
                           aggfunc=np.sum)
    table1 = pd.DataFrame(table.to_records())
#    table1 = table1.sort_values(by=[col1, value] , ascending=True)
    data = []
    for i in list(pd.unique(table1[col1])):
        label = list(table1[table1[col1]==i][col2])
        y_data = list(table1[table1[col1]==i][value])
        data.append(go.Bar(name=i, x=label, y=y_data,width=0.5))
    layout = go.Layout(barmode='stack',
                        xaxis = dict(title = '',titlefont=dict(size=15)),
                        yaxis = dict(title = 'Volume (in Truckload)',
                                     titlefont=dict(size=15)),width = 650,height = 350, margin=go.layout.Margin(l=50,r=1,b=100,t=1))
    fig5 = {'data': data, 'layout': layout}
#    plotly.io.orca.config.executable = 'C:/Rajkumar/dc_api/.env1/Scripts/plotly-orca-1.2.1-1/orca_app/orca.exe'
#    plotly.io.orca.config.save()
    pio.write_image(fig5, directory+image_dir+'fig5.png')
    vendor_image1 = directory+image_dir+'fig5.png'
    slide1.shapes.add_picture(vendor_image1,Inches(0.5),Inches(1.9))
    run_text(slide1,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
    return None

def col_converter(df, input1, output):
    if {input1}.issubset(df.columns):
        df.rename(columns={input1: output}, inplace=True)
    return None


def ppt(df):
    """Entry point of this module. It save ppt at location 
     ``FlaskWeb/Downloads/test-alpha.pptx``
    
    :param df: DataFrame of which ppt has to be made.
    :type df: DataFrame
    :param upto_number: Number upto which top should be calculated,
     defaults to 10
    :type upto_number: int, optional

    """
    df = df.apply(lambda x: x.str.strip().str.title() if(x.dtypes == object)
                  else x)
    col_converter(df, 'AnnualVolume', 'volume')
    col_converter(df, 'Rate', 'rate')
    col_converter(df, 'Spend', 'spend')
    col_converter(df, 'OriginName', v.Opoint)
    col_converter(df, 'DestinationName', v.Dpoint)
    col_converter(df, 'ReceivedDate', 'receiveddate')
    col_converter(df, 'TruckType', 'Trucktype')
    
    if {'spend', 'volume'}.issubset(df.columns):
        df['spend'] = pd.to_numeric(df['spend'], errors='coerce')
        df['volume'] = pd.to_numeric(df['volume'], errors='coerce')
    elif {'rate', 'volume'}.issubset(df.columns):
        df['rate'] = pd.to_numeric(df['rate'], errors='coerce')
        df['volume'] = pd.to_numeric(df['volume'], errors='coerce')
        df['spend'] = df['rate'] * df['volume']
    
    url = "https://logisticsnowtech3.blob.core.windows.net/rajkumar/sample.pptx"
    with urlopen(url) as f:
    	source = BytesIO(f.read())
    prs = Presentation(source) 
    Title_Only = 5
    try:
        summary_slide1 = prs.slides.add_slide(prs.slide_layouts[Title_Only])

        title_slide1 = summary_slide1.placeholders[0]
        title_slide1.text = "Dataset Summary"
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(2)
        shape = summary_slide1.shapes.add_table(8, 2, x, y, cx, cy)
        run_text(summary_slide1,4,6.70,10,1,"LogisticsNow Confidential",12,137,137,137)
        table = shape.table
        label = ['Parameter', 'Total Spend (in Lakh)', 'Total Volume',
                 'No. of Lanes', 'No. of Origin', 'No. of Destination',
                 'No. of Truck Type', 'No. of LSP']
        df_s1 = summarySlide(df)
        value = ['Value', df_s1['total spend'][0], df_s1['total vol'][0],
                 df_s1['# lane'][0], df_s1['# origin'][0],
                 df_s1['# destination'][0], df_s1['# unique_truck'][0],
                 df_s1['# unique_vendr'][0]]
        for idx in range(len(label)):
            cell = table.cell(idx, 0)
            cell.text = label[idx]
        for idx in range(len(value)):
            cell = table.cell(idx, 1)
            if isinstance(value[idx], str):
                cell.text = value[idx]
            elif isinstance(value[idx], (int, float, np.int64, np.float64)):
                cell.text = str(round(value[idx], 2))
                cell.alignment = PP_ALIGN.CENTER
    except Exception as e:
        print(str(e))
    
    try:
        Lane_map(prs, df, v.Opoint, v.Dpoint)
    except Exception as e:
        print(str(e))
    try:
        summary_chart(prs, df, 'receiveddate', 'volume')
    except Exception as e:
        print(str(e))
    try:
        bar_lineChart(prs, df, v.Opoint, 'spend', 'volume')
        bar_lineChart(prs, df, v.Dpoint, 'spend', 'volume')
    except Exception as e:
        print(str(e))
    # ----------------Lane-------------- #
    try:
        df['lane'] = df[v.Opoint] + ' -> ' + df[v.Dpoint]
        df_lane_vol = top(df, 'lane', "volume", 10)
        length = str(len(df_lane_vol))
        barChartConstruct(prs, title="Top {} Lanes by Volume".format(length),
                          typee='lane', df=df_lane_vol, based_on="volume",
                          axis_title="Volume (in Truckload)", color='g')
        df_lane_spend = top(df, 'lane', 'spend', 10, number_format=1e5)
        length = str(len(df_lane_spend))
        barChartConstruct(prs, title="Top {} Lanes by Spend".format(length),
                          typee='lane', df=df_lane_spend, based_on='spend',
                          axis_title="Spend (in Lakh)", color='b')
    except TypeError as e:
        print(e)
    try:
        donutChart(prs, df, 'Vendor', 'spend')
        donutChart(prs, df, 'Vendor', 'volume')
    except Exception as e:
        print(str(e))

    try:
        stackedBar(prs, df, 'Trucktype', 'Vendor', 'volume')
    except Exception as e:
        print(str(e))

    try:
        donutChart(prs, df, 'Trucktype', 'spend')
        donutChart(prs, df, 'Trucktype', 'volume')
    except Exception as e:
        print(str(e))
    move_slide(presentation=prs, old_index=1, new_index=13)
    prs.save(directory+image_dir+'test-alpha.pptx')
#    prs.save('FlaskWeb/Downloads/test-alpha.pptx')
    writeBlob(directory+image_dir+'test-alpha.pptx', 'test-alpha.pptx')


def main_func(request):
    req_body = request.get_json()
    file_url = req_body['file_url']
    if file_url is None:
        error = 'File not found'
        return f"{error}!"
    else:
        extension = os.path.splitext(file_url)[-1].lower()
    if extension == ".xlsx":
        df = pd.read_excel(file_url, encoding='utf-8')
    elif extension == ".csv":
        df = pd.read_csv(file_url, encoding='latin1')
    else:
        return f'please upload csv/xlsx file'
    try:
        ppt(df)
        return f'ppt created'
    except Exception as error:
        return f"{error}!"

